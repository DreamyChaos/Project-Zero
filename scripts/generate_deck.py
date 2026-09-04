import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def build_presentation():
    prs = Presentation()
    # 16:9 Widescreen aspect ratio
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Color Palette Definitions
    COLOR_BG = RGBColor(7, 11, 18)         # #070B12 Deep Space Black
    COLOR_CARD = RGBColor(15, 23, 42)      # #0F172A Subsurface Navy
    COLOR_BORDER = RGBColor(30, 41, 59)    # #1E293B Card Border
    COLOR_CYAN = RGBColor(0, 240, 255)      # #00F0FF Electric Cyan
    COLOR_GREEN = RGBColor(16, 185, 129)   # #10B981 Hydro Green
    COLOR_AMBER = RGBColor(245, 158, 11)   # #F59E0B Warning Amber
    COLOR_RED = RGBColor(239, 68, 68)      # #EF4444 Alert Red
    COLOR_TEXT_MAIN = RGBColor(248, 250, 252) # #F8FAFC
    COLOR_TEXT_MUTED = RGBColor(148, 163, 184) # #94A3B8
    COLOR_ACCENT_BLUE = RGBColor(2, 132, 199) # #0284C7

    IMG_DIR = r"C:\Users\bambi\.gemini\antigravity-ide\brain\2bb48f55-d8b7-4820-a3eb-37f08a573aab"
    COVER_IMG = os.path.join(IMG_DIR, "aquifer_cover_1786908412111.jpg")
    GIS_IMG = os.path.join(IMG_DIR, "satellite_gis_grid_1786908427998.jpg")

    def add_bg(slide):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = COLOR_BG
        bg.line.fill.background()
        return bg

    def add_header(slide, num_str, title_str, subtitle_str=""):
        # Category Tag
        tx_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.4))
        tf = tx_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"{num_str}  // GROUNDWATER INTELLIGENCE PLATFORM"
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = COLOR_CYAN
        p.font.name = 'Segoe UI'

        # Main Slide Title
        tx_box2 = slide.shapes.add_textbox(Inches(0.8), Inches(0.75), Inches(11.7), Inches(0.7))
        tf2 = tx_box2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = title_str
        p2.font.size = Pt(24)
        p2.font.bold = True
        p2.font.color.rgb = COLOR_TEXT_MAIN
        p2.font.name = 'Segoe UI'

        # Subtitle if present
        if subtitle_str:
            p3 = tf2.add_paragraph()
            p3.text = subtitle_str
            p3.font.size = Pt(13)
            p3.font.color.rgb = COLOR_TEXT_MUTED
            p3.font.name = 'Segoe UI'

    def create_card(slide, left, top, width, height, title="", border_color=COLOR_BORDER):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_CARD
        card.line.color.rgb = border_color
        card.line.width = Pt(1.5)
        
        if title:
            tb = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.15), width - Inches(0.4), Inches(0.4))
            tf = tb.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = title.upper()
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = COLOR_CYAN
            p.font.name = 'Segoe UI'
        return card

    # ==================== SLIDE 01: COVER ====================
    slide1 = prs.slides.add_slide(blank_layout)
    add_bg(slide1)
    if os.path.exists(COVER_IMG):
        slide1.shapes.add_picture(COVER_IMG, Inches(4.5), Inches(0.5), Inches(8.3), Inches(6.5))
    
    # Overlay card for Title
    card1 = create_card(slide1, Inches(0.8), Inches(1.5), Inches(6.5), Inches(4.8), border_color=COLOR_CYAN)
    tb = slide1.shapes.add_textbox(Inches(1.1), Inches(1.8), Inches(5.9), Inches(4.2))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "GROUNDWATER INTELLIGENCE PLATFORM"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT_MAIN

    p2 = tf.add_paragraph()
    p2.text = "\nA Data-Driven System for Groundwater Monitoring, Spatial Mapping & Predictive Decision Support"
    p2.font.size = Pt(15)
    p2.font.color.rgb = COLOR_CYAN

    p3 = tf.add_paragraph()
    p3.text = "\nCore Narrative Architecture:\nSENSE → INTEGRATE → MAP → ANALYZE → PREDICT → DECIDE"
    p3.font.size = Pt(12)
    p3.font.color.rgb = COLOR_TEXT_MUTED

    p4 = tf.add_paragraph()
    p4.text = "\nPowered by Hardware Telemetry, GIS, Remote Sensing & AI/ML"
    p4.font.size = Pt(11)
    p4.font.color.rgb = COLOR_GREEN

    # ==================== SLIDE 02: THE PROBLEM ====================
    slide2 = prs.slides.add_slide(blank_layout)
    add_bg(slide2)
    add_header(slide2, "02", "GROUNDWATER IS INVISIBLE — ITS IMPACT ISN'T.", "Critical global challenges facing groundwater monitoring and management")

    col_w = Inches(3.6)
    c1 = create_card(slide2, Inches(0.8), Inches(1.8), col_w, Inches(5.0), title="01 / Aquifer Depletion", border_color=COLOR_RED)
    tb1 = slide2.shapes.add_textbox(Inches(1.0), Inches(2.3), col_w - Inches(0.4), Inches(4.2))
    tf1 = tb1.text_frame
    tf1.word_wrap = True
    tf1.paragraphs[0].text = "Unseen Reserve Over-extraction"
    tf1.paragraphs[0].font.bold = True
    tf1.paragraphs[0].font.color.rgb = COLOR_TEXT_MAIN
    tf1.paragraphs[0].font.size = Pt(14)
    p = tf1.add_paragraph()
    p.text = "\n• Rapid decline in water table levels across agricultural & urban zones.\n• Severe lag in traditional manual well measurements.\n• Subsidence & irreversible storage capacity destruction."
    p.font.size = Pt(12)
    p.font.color.rgb = COLOR_TEXT_MUTED

    c2 = create_card(slide2, Inches(4.8), Inches(1.8), col_w, Inches(5.0), title="02 / Water Quality Risk", border_color=COLOR_AMBER)
    tb2 = slide2.shapes.add_textbox(Inches(5.0), Inches(2.3), col_w - Inches(0.4), Inches(4.2))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    tf2.paragraphs[0].text = "Silent Contamination & Salinization"
    tf2.paragraphs[0].font.bold = True
    tf2.paragraphs[0].font.color.rgb = COLOR_TEXT_MAIN
    tf2.paragraphs[0].font.size = Pt(14)
    p = tf2.add_paragraph()
    p.text = "\n• Rising Electrical Conductivity (EC) & Turbidity levels.\n• pH spikes endangering drinking water & crop yield.\n• Lack of continuous continuous temporal contamination tracking."
    p.font.size = Pt(12)
    p.font.color.rgb = COLOR_TEXT_MUTED

    c3 = create_card(slide2, Inches(8.8), Inches(1.8), col_w, Inches(5.0), title="03 / Blind Decision Making", border_color=COLOR_AMBER)
    tb3 = slide2.shapes.add_textbox(Inches(9.0), Inches(2.3), col_w - Inches(0.4), Inches(4.2))
    tf3 = tb3.text_frame
    tf3.word_wrap = True
    tf3.paragraphs[0].text = "Reactive Management Failure"
    tf3.paragraphs[0].font.bold = True
    tf3.paragraphs[0].font.color.rgb = COLOR_TEXT_MAIN
    tf3.paragraphs[0].font.size = Pt(14)
    p = tf3.add_paragraph()
    p.text = "\n• Policy interventions arrive after wells run dry.\n• Inability to model future stress under climate variability.\n• Absence of centralized spatial intelligence."
    p.font.size = Pt(12)
    p.font.color.rgb = COLOR_TEXT_MUTED

    # ==================== SLIDE 03: THE DATA GAP ====================
    slide3 = prs.slides.add_slide(blank_layout)
    add_bg(slide3)
    add_header(slide3, "03", "THE DATA EXISTS. THE CONNECTION DOESN'T.", "Groundwater data is trapped in isolated, incompatible silos")

    silo_w = Inches(2.2)
    silos = [
        ("Physical Sensors", "Local well depth & water metrics", COLOR_CYAN),
        ("Satellite / Remote", "Soil moisture & GRACE gravity data", COLOR_ACCENT_BLUE),
        ("Rainfall Data", "Meteorological precipitation logs", COLOR_GREEN),
        ("Soil & Terrain", "Infiltration & geological GIS layers", COLOR_AMBER),
        ("Historical Logs", "Legacy paper records & state reports", COLOR_TEXT_MUTED)
    ]
    for idx, (stitle, sdesc, scolor) in enumerate(silos):
        left_pos = Inches(0.8 + idx * 2.4)
        create_card(slide3, left_pos, Inches(2.0), silo_w, Inches(3.2), title=f"Silo 0{idx+1}", border_color=scolor)
        tb = slide3.shapes.add_textbox(left_pos + Inches(0.15), Inches(2.6), silo_w - Inches(0.3), Inches(2.4))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].text = stitle
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.size = Pt(13)
        tf.paragraphs[0].font.color.rgb = COLOR_TEXT_MAIN
        p = tf.add_paragraph()
        p.text = f"\n{sdesc}\n\n[ Isolated ]"
        p.font.size = Pt(11)
        p.font.color.rgb = COLOR_TEXT_MUTED

    # Impact Card Bottom
    create_card(slide3, Inches(0.8), Inches(5.5), Inches(11.7), Inches(1.3), title="THE CONSEQUENCE", border_color=COLOR_RED)
    tb_bot = slide3.shapes.add_textbox(Inches(1.0), Inches(5.85), Inches(11.3), Inches(0.8))
    tf_bot = tb_bot.text_frame
    tf_bot.word_wrap = True
    p = tf_bot.paragraphs[0]
    p.text = "Without unification, decision-makers get partial views, delayed warnings, and fragmented policy actions."
    p.font.size = Pt(13)
    p.font.color.rgb = COLOR_TEXT_MAIN

    # ==================== SLIDE 04: OUR VISION ====================
    slide4 = prs.slides.add_slide(blank_layout)
    add_bg(slide4)
    add_header(slide4, "04", "ONE PLATFORM. MULTIPLE DATA SOURCES. ONE GROUNDWATER VIEW.", "Fusing physical telemetry with spatial environmental intelligence")

    create_card(slide4, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.9), title="UNIFIED PLATFORM ARCHITECTURE", border_color=COLOR_CYAN)
    tb = slide4.shapes.add_textbox(Inches(1.0), Inches(2.3), Inches(5.2), Inches(4.2))
    tf = tb.text_frame
    tf.word_wrap = True
    bullets = [
        ("Ground-Truth Telemetry", "Real-time observations from physical IoT nodes (Raspberry Pi 3B+ & ESP32)."),
        ("Remote Sensing Integration", "Broad spatial coverage extending context beyond physical sensor sites."),
        ("GIS Layer Fusion", "Topography, slope, land cover, and soil characteristics mapped seamlessly."),
        ("Centralized Data Engine", "Single source of truth feeding spatial analytics and predictive intelligence.")
    ]
    for btitle, bdesc in bullets:
        p = tf.add_paragraph()
        p.text = f"• {btitle}: "
        p.font.bold = True
        p.font.size = Pt(12)
        p.font.color.rgb = COLOR_CYAN
        p_sub = tf.add_paragraph()
        p_sub.text = f"  {bdesc}\n"
        p_sub.font.size = Pt(11)
        p_sub.font.color.rgb = COLOR_TEXT_MUTED

    # Convergence Diagram right box
    create_card(slide4, Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.9), title="THE CORE VALUE NARRATIVE", border_color=COLOR_GREEN)
    tb_right = slide4.shapes.add_textbox(Inches(7.1), Inches(2.4), Inches(5.1), Inches(4.0))
    tf_r = tb_right.text_frame
    tf_r.word_wrap = True
    steps = [
        "1. SENSE  → Real-time physical sensor data collection",
        "2. INTEGRATE → Pipeline unification of multi-source datasets",
        "3. MAP → GIS spatial potential & risk visualization",
        "4. ANALYZE → Time-series depletion & quality trend evaluation",
        "5. PREDICT → AI/ML predictive forecasting",
        "6. DECIDE → Government decision-support dashboard"
    ]
    for s in steps:
        p = tf_r.add_paragraph()
        p.text = s
        p.font.bold = True
        p.font.size = Pt(13)
        p.font.color.rgb = COLOR_TEXT_MAIN
        tf_r.add_paragraph().text = ""

    # ==================== SLIDE 05: SYSTEM ARCHITECTURE ====================
    slide5 = prs.slides.add_slide(blank_layout)
    add_bg(slide5)
    add_header(slide5, "05", "SYSTEM ARCHITECTURE", "End-to-end data pipeline from physical sensing to decision support")

    layers = [
        ("01 / PHYSICAL LAYER", "Raspberry Pi 3B+ | ESP32 | 6 Telemetry Sensors | Custom PCB | Physical Model", COLOR_CYAN),
        ("02 / GATEWAY & API LAYER", "Data Acquisition Gateway | Hardware Data API | Validation & Cleaning", COLOR_ACCENT_BLUE),
        ("03 / DATA BASE & INTEGRATION", "Central Database | GIS Spatial Layer Index | Remote Sensing Ingestion", COLOR_GREEN),
        ("04 / ANALYTICS & MAP ENGINE", "Spatial Potential Mapping | Recharge Calculation | Quality Risk Indexing", COLOR_AMBER),
        ("05 / PREDICTIVE INTELLIGENCE", "AI/ML Prediction Engine | Trend Analytics | Stress Forecasting", COLOR_RED),
        ("06 / DECISION PLATFORM", "Web Frontend | GIS Dashboard | Government Decision Support & Reports", COLOR_TEXT_MAIN)
    ]
    for idx, (ltitle, ldesc, lcolor) in enumerate(layers):
        top_pos = Inches(1.8 + idx * 0.85)
        card = create_card(slide5, Inches(0.8), top_pos, Inches(11.7), Inches(0.75), border_color=lcolor)
        tb = slide5.shapes.add_textbox(Inches(1.0), top_pos + Inches(0.12), Inches(11.3), Inches(0.5))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"{ltitle}  "
        p.font.bold = True
        p.font.size = Pt(12)
        p.font.color.rgb = lcolor
        
        p2 = tf.add_paragraph()
        p2.text = f"    {ldesc}"
        p2.font.size = Pt(11)
        p2.font.color.rgb = COLOR_TEXT_MUTED

    # ==================== SLIDE 06: FIELD DATA NODE ====================
    slide6 = prs.slides.add_slide(blank_layout)
    add_bg(slide6)
    add_header(slide6, "06", "FIELD DATA NODE", "Physical prototype data-acquisition and edge gateway architecture")

    # Left Hardware Card
    create_card(slide6, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.9), title="HARDWARE CORE COMPONENTS", border_color=COLOR_CYAN)
    tb_hw = slide6.shapes.add_textbox(Inches(1.0), Inches(2.3), Inches(5.2), Inches(4.2))
    tf_hw = tb_hw.text_frame
    tf_hw.word_wrap = True
    hw_list = [
        ("Raspberry Pi 3B+", "Central edge gateway, local buffering & API communication"),
        ("ESP32 Microcontroller", "Low-latency ADC reading & sensor sampling controller"),
        ("Custom PCB & Wiring", "Noise reduction, regulated power delivery & connector bus"),
        ("Pump + Tubing", "Controlled fluid circulation for dynamic sensor testing"),
        ("Physical Groundwater Model", "Simulated aquifer column for continuous physical validation")
    ]
    for htitle, hdesc in hw_list:
        p = tf_hw.add_paragraph()
        p.text = f"• {htitle}: "
        p.font.bold = True
        p.font.size = Pt(11)
        p.font.color.rgb = COLOR_TEXT_MAIN
        p_sub = tf_hw.add_paragraph()
        p_sub.text = f"  {hdesc}"
        p_sub.font.size = Pt(10)
        p_sub.font.color.rgb = COLOR_TEXT_MUTED

    # Right Architecture Box
    create_card(slide6, Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.9), title="EDGE GATEWAY ROLE IN SYSTEM", border_color=COLOR_GREEN)
    tb_role = slide6.shapes.add_textbox(Inches(7.1), Inches(2.4), Inches(5.1), Inches(4.0))
    tf_role = tb_role.text_frame
    tf_role.word_wrap = True
    p = tf_role.paragraphs[0]
    p.text = "Physical Gateway vs Platform Positioning:"
    p.font.bold = True
    p.font.size = Pt(14)
    p.font.color.rgb = COLOR_CYAN

    roles = [
        ("Data Ground-Truth", "Provides real-world empirical validation for spatial models."),
        ("Edge Data Acquisition", "Continuous, autonomous sampling of water depth & quality parameters."),
        ("Non-Exclusive Role", "The hardware node is an essential data source, but the software intelligence platform is the primary project engine.")
    ]
    for rtitle, rdesc in roles:
        p = tf_role.add_paragraph()
        p.text = f"\n✓ {rtitle}"
        p.font.bold = True
        p.font.size = Pt(12)
        p.font.color.rgb = COLOR_GREEN
        p2 = tf_role.add_paragraph()
        p2.text = f"   {rdesc}"
        p2.font.size = Pt(11)
        p2.font.color.rgb = COLOR_TEXT_MUTED

    # ==================== SLIDE 07: SENSOR DATA ====================
    slide7 = prs.slides.add_slide(blank_layout)
    add_bg(slide7)
    add_header(slide7, "07", "SENSOR INSTRUMENTATION & TELEMETRY", "Key physical parameters collected by the edge sensor payload")

    sensors = [
        ("Groundwater Level", "Ultrasonic / Pressure depth sensor", "Tracks real-time water table height & drawdown", COLOR_CYAN),
        ("pH Sensor", "Electrochemical electrode", "Monitors acidity/alkalinity for water safety", COLOR_GREEN),
        ("Electrical Conductivity (EC)", "Conductivity probe", "Detects total dissolved solids & salinization", COLOR_AMBER),
        ("Turbidity Sensor", "Optical scattering probe", "Measures suspended solids & sedimentation", COLOR_RED),
        ("Temperature Sensor", "Digital thermal probe", "Measures water/ambient thermal variations", COLOR_ACCENT_BLUE),
        ("Soil-Moisture Sensor", "Capacitive moisture probe", "Evaluates surface infiltration readiness", COLOR_TEXT_MAIN)
    ]

    for idx, (stitle, ssub, sdesc, scolor) in enumerate(sensors):
        col = idx % 3
        row = idx // 3
        left_pos = Inches(0.8 + col * 4.0)
        top_pos = Inches(1.8 + row * 2.5)

        create_card(slide7, left_pos, top_pos, Inches(3.7), Inches(2.2), title=stitle, border_color=scolor)
        tb = slide7.shapes.add_textbox(left_pos + Inches(0.15), top_pos + Inches(0.55), Inches(3.4), Inches(1.5))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].text = ssub
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.size = Pt(11)
        tf.paragraphs[0].font.color.rgb = COLOR_TEXT_MAIN

        p = tf.add_paragraph()
        p.text = f"\n{sdesc}"
        p.font.size = Pt(10)
        p.font.color.rgb = COLOR_TEXT_MUTED

    # ==================== SLIDE 08: GIS INTERACTIVE MAP ====================
    slide8 = prs.slides.add_slide(blank_layout)
    add_bg(slide8)
    add_header(slide8, "08", "GIS INTERACTIVE MAP", "Spatial intelligence layer providing multi-resolution geographic visualization")

    create_card(slide8, Inches(0.8), Inches(1.8), Inches(6.0), Inches(4.9), title="SPATIAL MAPPING CAPABILITIES", border_color=COLOR_CYAN)
    tb = slide8.shapes.add_textbox(Inches(1.0), Inches(2.3), Inches(5.6), Inches(4.2))
    tf = tb.text_frame
    tf.word_wrap = True
    gis_features = [
        ("Multi-Layer Geospatial Overlay", "Integrate sensor nodes, administrative boundaries, terrain contours, and watershed basins."),
        ("Interactive Coordinate Querying", "Clickable spatial nodes displaying live telemetry, depth history, and risk scores."),
        ("Raster & Vector Processing", "Smooth rendering of satellite rasters alongside point sensor locations."),
        ("Dynamic Heatmap Rendering", "Interpolate spatial data across unmonitored zones to estimate groundwater potential.")
    ]
    for ftitle, fdesc in gis_features:
        p = tf.add_paragraph()
        p.text = f"• {ftitle}: "
        p.font.bold = True
        p.font.size = Pt(11)
        p.font.color.rgb = COLOR_CYAN
        p_sub = tf.add_paragraph()
        p_sub.text = f"  {fdesc}\n"
        p_sub.font.size = Pt(10)
        p_sub.font.color.rgb = COLOR_TEXT_MUTED

    # Right side: Satellite GIS image background card
    if os.path.exists(GIS_IMG):
        slide8.shapes.add_picture(GIS_IMG, Inches(7.1), Inches(1.8), Inches(5.4), Inches(4.9))

    # ==================== SLIDE 09: REMOTE SENSING ====================
    slide9 = prs.slides.add_slide(blank_layout)
    add_bg(slide9)
    add_header(slide9, "09", "REMOTE SENSING & ENVIRONMENTAL FUSION", "Expanding coverage beyond point-sensor locations")

    cards_rs = [
        ("01 / Satellite Imagery", "Multi-spectral optical & radar observation for vegetation index (NDVI) & surface moisture.", COLOR_CYAN),
        ("02 / Environmental Datasets", "Integration of precipitation, temperature, evapotranspiration, and humidity logs.", COLOR_GREEN),
        ("03 / Macro Spatial Scale", "Connecting localized well sensors to regional aquifer-scale dynamics.", COLOR_AMBER)
    ]

    for idx, (title, desc, color) in enumerate(cards_rs):
        left = Inches(0.8 + idx * 4.0)
        create_card(slide9, left, Inches(1.8), Inches(3.7), Inches(4.9), title=title, border_color=color)
        tb = slide9.shapes.add_textbox(left + Inches(0.2), Inches(2.4), Inches(3.3), Inches(4.0))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(12)
        p.font.color.rgb = COLOR_TEXT_MAIN

        p2 = tf.add_paragraph()
        p2.text = "\n\nKey Benefit:"
        p2.font.bold = True
        p2.font.color.rgb = color
        p2.font.size = Pt(11)

        p3 = tf.add_paragraph()
        p3.text = "Fills spatial gaps between physical monitoring wells without requiring thousands of physical sensors."
        p3.font.size = Pt(10)
        p3.font.color.rgb = COLOR_TEXT_MUTED

    # ==================== SLIDE 10: GROUNDWATER POTENTIAL MAPPING ====================
    slide10 = prs.slides.add_slide(blank_layout)
    add_bg(slide10)
    add_header(slide10, "10", "GROUNDWATER POTENTIAL MAPPING", "Multi-criteria spatial assessment for groundwater availability")

    create_card(slide10, Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0), title="MULTI-LAYER OVERLAY MATRIX", border_color=COLOR_CYAN)

    tb = slide10.shapes.add_textbox(Inches(1.1), Inches(2.4), Inches(11.1), Inches(4.2))
    tf = tb.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "Groundwater Potential Assessment Framework (Illustrative Model)"
    p.font.bold = True
    p.font.size = Pt(14)
    p.font.color.rgb = COLOR_TEXT_MAIN

    p2 = tf.add_paragraph()
    p2.text = "\nInputs: Rainfall + Soil Type + Topography/Slope + Land Cover + Drainage Density + Sensor Ground Truth\n"
    p2.font.size = Pt(12)
    p2.font.color.rgb = COLOR_CYAN

    factors = [
        ("Topography & Slope", "Flatter terrain promotes water ponding and infiltration over rapid runoff."),
        ("Soil Permeability", "Sandy & porous soils enable higher percolation rates than dense clay."),
        ("Drainage Density", "Lower stream density indicates higher subsurface absorption capacity."),
        ("Land Cover & NDVI", "Dense vegetation enhances root-zone soil porosity and reduces direct evaporation.")
    ]

    for ftitle, fdesc in factors:
        p = tf.add_paragraph()
        p.text = f"• {ftitle}: "
        p.font.bold = True
        p.font.color.rgb = COLOR_TEXT_MAIN
        p.font.size = Pt(11)
        p2 = tf.add_paragraph()
        p2.text = f"  {fdesc}"
        p2.font.size = Pt(10)
        p2.font.color.rgb = COLOR_TEXT_MUTED

    # ==================== SLIDE 11: RECHARGE-POTENTIAL MAPPING ====================
    slide11 = prs.slides.add_slide(blank_layout)
    add_bg(slide11)
    add_header(slide11, "11", "RECHARGE-POTENTIAL MAPPING", "Identifying optimal zones for artificial & natural groundwater recharge")

    col_w = Inches(5.6)
    create_card(slide11, Inches(0.8), Inches(1.8), col_w, Inches(4.9), title="INFILTRATION & RECHARGE FACTORS", border_color=COLOR_GREEN)
    tb1 = slide11.shapes.add_textbox(Inches(1.0), Inches(2.3), col_w - Inches(0.4), Inches(4.2))
    tf1 = tb1.text_frame
    tf1.word_wrap = True
    tf1.paragraphs[0].text = "Spatial Recharge Factors:"
    tf1.paragraphs[0].font.bold = True
    tf1.paragraphs[0].font.color.rgb = COLOR_TEXT_MAIN
    tf1.paragraphs[0].font.size = Pt(13)

    r_factors = [
        "Precipitation Intensity & Duration",
        "Soil Hydraulic Conductivity",
        "Geological Lithology & Subsurface Permeability",
        "Surface Slope & Runoff Velocity",
        "Soil Moisture Saturation Levels"
    ]
    for rf in r_factors:
        p = tf1.add_paragraph()
        p.text = f"• {rf}"
        p.font.size = Pt(11)
        p.font.color.rgb = COLOR_TEXT_MUTED

    create_card(slide11, Inches(6.8), Inches(1.8), col_w, Inches(4.9), title="DECISION SUPPORT APPLICATIONS", border_color=COLOR_CYAN)
    tb2 = slide11.shapes.add_textbox(Inches(7.0), Inches(2.3), col_w - Inches(0.4), Inches(4.2))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    tf2.paragraphs[0].text = "Actionable Interventions:"
    tf2.paragraphs[0].font.bold = True
    tf2.paragraphs[0].font.color.rgb = COLOR_TEXT_MAIN
    tf2.paragraphs[0].font.size = Pt(13)

    actions = [
        ("Check Dam Site Selection", "Identify high-percolation zones for check dams."),
        ("Recharge Shaft Placement", "Locate optimal sites for rainwater injection shafts."),
        ("Rainwater Harvesting Strategy", "Guide municipal rainwater collection mandates."),
        ("Conservation Zoning", "Designate protected high-recharge watershed basins.")
    ]
    for atitle, adesc in actions:
        p = tf2.add_paragraph()
        p.text = f"✓ {atitle}: "
        p.font.bold = True
        p.font.size = Pt(11)
        p.font.color.rgb = COLOR_GREEN
        p_sub = tf2.add_paragraph()
        p_sub.text = f"  {adesc}"
        p_sub.font.size = Pt(10)
        p_sub.font.color.rgb = COLOR_TEXT_MUTED

    # ==================== SLIDE 12: WATER-QUALITY RISK MAPPING ====================
    slide12 = prs.slides.add_slide(blank_layout)
    add_bg(slide12)
    add_header(slide12, "12", "WATER-QUALITY RISK MAPPING", "Transforming point sensor measurements into spatial risk maps")

    q_cards = [
        ("pH Spatial Index", "Acidity & Alkalinity", "Maps agricultural runoff, chemical leaching & industrial contamination zones.", COLOR_CYAN),
        ("EC Risk Index", "Electrical Conductivity", "Tracks salinization, mineral concentration & coastal saltwater intrusion.", COLOR_AMBER),
        ("Turbidity Risk Index", "Suspended Sediment", "Identifies wellhead structural damage, sediment ingress & filtration failure.", COLOR_RED)
    ]
    for idx, (qtitle, qsub, qdesc, qcolor) in enumerate(q_cards):
        left = Inches(0.8 + idx * 4.0)
        create_card(slide12, left, Inches(1.8), Inches(3.7), Inches(3.2), title=qtitle, border_color=qcolor)
        tb = slide12.shapes.add_textbox(left + Inches(0.2), Inches(2.4), Inches(3.3), Inches(2.4))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].text = qsub
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.size = Pt(12)
        tf.paragraphs[0].font.color.rgb = COLOR_TEXT_MAIN
        p = tf.add_paragraph()
        p.text = f"\n{qdesc}"
        p.font.size = Pt(10)
        p.font.color.rgb = COLOR_TEXT_MUTED

    # Bottom Integrated Index Card
    create_card(slide12, Inches(0.8), Inches(5.2), Inches(11.7), Inches(1.6), title="INTEGRATED WATER-QUALITY RISK SCORE", border_color=COLOR_GREEN)
    tb_bot = slide12.shapes.add_textbox(Inches(1.0), Inches(5.55), Inches(11.3), Inches(1.1))
    tf_bot = tb_bot.text_frame
    tf_bot.word_wrap = True
    p = tf_bot.paragraphs[0]
    p.text = "Spatial interpolation combines pH, EC, and Turbidity into a unified Risk Index (Safe / Moderate / Critical Zone) to alert authorities prior to public water supply contamination."
    p.font.size = Pt(12)
    p.font.color.rgb = COLOR_TEXT_MAIN

    # ==================== SLIDE 13: HISTORICAL TREND ANALYSIS ====================
    slide13 = prs.slides.add_slide(blank_layout)
    add_bg(slide13)
    add_header(slide13, "13", "HISTORICAL TREND ANALYSIS", "Time-series visualization of water table depth and environmental factors")

    create_card(slide13, Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0), title="TEMPORAL PATTERN EVALUATION", border_color=COLOR_CYAN)
    tb = slide13.shapes.add_textbox(Inches(1.1), Inches(2.4), Inches(11.1), Inches(4.2))
    tf = tb.text_frame
    tf.word_wrap = True

    tf.paragraphs[0].text = "Key Analytical Dimensions:"
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.color.rgb = COLOR_TEXT_MAIN

    trends = [
        ("Seasonal Water-Table Fluctuations", "Comparing monsoon recharge peaks against summer drawdown troughs."),
        ("Multi-Year Depletion Trajectories", "Identifying long-term negative baseline drift across consecutive years."),
        ("Rainfall-Recharge Lag Analysis", "Quantifying the precise response time between rainfall events and water table rise."),
        ("Anomalous Event Detection", "Flagging sudden unseasonal drops caused by illegal over-pumping or sensor fault.")
    ]
    for ttitle, tdesc in trends:
        p = tf.add_paragraph()
        p.text = f"\n• {ttitle}: "
        p.font.bold = True
        p.font.color.rgb = COLOR_CYAN
        p.font.size = Pt(12)
        p_sub = tf.add_paragraph()
        p_sub.text = f"  {tdesc}"
        p_sub.font.size = Pt(11)
        p_sub.font.color.rgb = COLOR_TEXT_MUTED

    # ==================== SLIDE 14: GROUNDWATER DEPLETION / STRESS ====================
    slide14 = prs.slides.add_slide(blank_layout)
    add_bg(slide14)
    add_header(slide14, "14", "GROUNDWATER DEPLETION & STRESS ZONES", "Identifying critical zones approaching threshold depletion")

    create_card(slide14, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.9), title="STRESS CLASSIFICATION MATRIX", border_color=COLOR_AMBER)
    tb1 = slide14.shapes.add_textbox(Inches(1.0), Inches(2.3), Inches(5.2), Inches(4.2))
    tf1 = tb1.text_frame
    tf1.word_wrap = True

    levels = [
        ("Safe Zone", "Stable water table, positive annual recharge balance.", COLOR_GREEN),
        ("Semi-Critical Zone", "Minor drawdown trend, extraction equals natural recharge.", COLOR_AMBER),
        ("Critical Stress Zone", "Severe persistent drawdown, extraction exceeds recharge by >20%.", COLOR_RED),
        ("Over-Exploited Zone", "Irreversible storage decline, risk of well failure & land subsidence.", COLOR_RED)
    ]
    for ltitle, ldesc, lcolor in levels:
        p = tf1.add_paragraph()
        p.text = f"■ {ltitle}: "
        p.font.bold = True
        p.font.color.rgb = lcolor
        p.font.size = Pt(11)
        p_sub = tf1.add_paragraph()
        p_sub.text = f"  {ldesc}\n"
        p_sub.font.size = Pt(10)
        p_sub.font.color.rgb = COLOR_TEXT_MUTED

    create_card(slide14, Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.9), title="EARLY WARNING METRICS", border_color=COLOR_RED)
    tb2 = slide14.shapes.add_textbox(Inches(7.0), Inches(2.3), Inches(5.3), Inches(4.2))
    tf2 = tb2.text_frame
    tf2.word_wrap = True

    metrics = [
        ("Rate of Drawdown (m/year)", "Accelerating annual depletion rate."),
        ("Recharge Deficit Score", "Percentage of unreplenished extracted volume."),
        ("Water Quality Degradation Rate", "Simultaneous rise in EC/salinity alongside depth drop."),
        ("Proximity to Critical Infrastructure", "Threat level to municipal supply wells.")
    ]
    for mtitle, mdesc in metrics:
        p = tf2.add_paragraph()
        p.text = f"⚠ {mtitle}: "
        p.font.bold = True
        p.font.color.rgb = COLOR_TEXT_MAIN
        p.font.size = Pt(11)
        p_sub = tf2.add_paragraph()
        p_sub.text = f"  {mdesc}\n"
        p_sub.font.size = Pt(10)
        p_sub.font.color.rgb = COLOR_TEXT_MUTED

    # ==================== SLIDE 15: AI / ML ENGINE ====================
    slide15 = prs.slides.add_slide(blank_layout)
    add_bg(slide15)
    add_header(slide15, "15", "AI / ML PREDICTION ENGINE", "Predictive Intelligence Layer for groundwater forecasting (Algorithm Agnostic)")

    # Top banner emphasizing requirement 3
    create_card(slide15, Inches(0.8), Inches(1.8), Inches(11.7), Inches(1.0), title="PREDICTIVE INTELLIGENCE ARCHITECTURE", border_color=COLOR_CYAN)
    tb_top = slide15.shapes.add_textbox(Inches(1.0), Inches(2.15), Inches(11.3), Inches(0.6))
    tf_top = tb_top.text_frame
    tf_top.word_wrap = True
    p = tf_top.paragraphs[0]
    p.text = "Modular AI/ML engine designed to evaluate temporal patterns and environmental features. Model selection (e.g., ensemble, neural, or time-series regression) will be finalized based on dataset resolution and validation performance."
    p.font.size = Pt(11)
    p.font.color.rgb = COLOR_TEXT_MAIN

    # Left Box: Data Inputs
    create_card(slide15, Inches(0.8), Inches(3.0), Inches(5.6), Inches(3.7), title="FEATURE PIPELINE INPUTS", border_color=COLOR_ACCENT_BLUE)
    tb_in = slide15.shapes.add_textbox(Inches(1.0), Inches(3.5), Inches(5.2), Inches(3.0))
    tf_in = tb_in.text_frame
    tf_in.word_wrap = True
    in_features = [
        "Historical Water Level Time Series",
        "Precipitation & Evapotranspiration Rates",
        "Soil Moisture & Satellite NDVI Rasters",
        "Geological Lithology & Infiltration Rates",
        "Live Edge Sensor Telemetry Feeds"
    ]
    for inf in in_features:
        p = tf_in.add_paragraph()
        p.text = f"→ {inf}"
        p.font.size = Pt(11)
        p.font.color.rgb = COLOR_TEXT_MUTED

    # Right Box: Target Outputs
    create_card(slide15, Inches(6.8), Inches(3.0), Inches(5.7), Inches(3.7), title="PREDICTIVE OUTPUTS", border_color=COLOR_GREEN)
    tb_out = slide15.shapes.add_textbox(Inches(7.0), Inches(3.5), Inches(5.3), Inches(3.0))
    tf_out = tb_out.text_frame
    tf_out.word_wrap = True
    out_features = [
        "30 / 60 / 90-Day Water Level Forecasts",
        "Depletion Risk Probability Distribution",
        "Recharge Deficit Predictions",
        "Water Quality Risk Trend Horizon",
        "Uncertainty Bounds & Confidence Intervals"
    ]
    for outf in out_features:
        p = tf_out.add_paragraph()
        p.text = f"★ {outf}"
        p.font.size = Pt(11)
        p.font.color.rgb = COLOR_TEXT_MAIN

    # ==================== SLIDE 16: FUTURE STRESS PREDICTION ====================
    slide16 = prs.slides.add_slide(blank_layout)
    add_bg(slide16)
    add_header(slide16, "16", "FUTURE GROUNDWATER-STRESS PREDICTION", "Scenario-based forecasting for climate variability and extraction demand")

    scenarios = [
        ("Baseline Scenario", "Normal Rainfall & Current Extraction", "Projects baseline water table behavior under standard seasonal precipitation.", COLOR_GREEN),
        ("Monsoon Deficit Scenario", "Below-Average Rainfall (-30%)", "Forecasts severe drawdown speed during drought years to prepare buffer storage.", COLOR_AMBER),
        ("High Extraction Demand", "Agricultural/Urban Expansion (+25%)", "Simulates aquifer stress under accelerated pumping schedules.", COLOR_RED)
    ]

    for idx, (stitle, ssub, sdesc, scolor) in enumerate(scenarios):
        left = Inches(0.8 + idx * 4.0)
        create_card(slide16, left, Inches(1.8), Inches(3.7), Inches(4.9), title=stitle, border_color=scolor)
        tb = slide16.shapes.add_textbox(left + Inches(0.2), Inches(2.4), Inches(3.3), Inches(4.0))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].text = ssub
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.size = Pt(12)
        tf.paragraphs[0].font.color.rgb = COLOR_TEXT_MAIN
        p = tf.add_paragraph()
        p.text = f"\n{sdesc}"
        p.font.size = Pt(11)
        p.font.color.rgb = COLOR_TEXT_MUTED

    # ==================== SLIDE 17: DECISION-SUPPORT DASHBOARD ====================
    slide17 = prs.slides.add_slide(blank_layout)
    add_bg(slide17)
    add_header(slide17, "17", "DECISION-SUPPORT DASHBOARD", "Turning complex spatial data and predictions into government policy action")

    dash_widgets = [
        ("01 / Groundwater Status Overview", "Regional water table health, active well count, live telemetry feeds.", COLOR_CYAN),
        ("02 / Recharge Priority Map", "Geospatial highlighting of top-priority zones for check dams & artificial recharge.", COLOR_GREEN),
        ("03 / Contamination Risk Alerts", "Automated warnings for regions exceeding safe pH, EC, or turbidity thresholds.", COLOR_AMBER),
        ("04 / Depletion Trend Forecasting", "Interactive slider to visualize projected groundwater stress over 1–5 years.", COLOR_RED)
    ]

    for idx, (wtitle, wdesc, wcolor) in enumerate(dash_widgets):
        col = idx % 2
        row = idx // 2
        left = Inches(0.8 + col * 5.9)
        top = Inches(1.8 + row * 2.5)

        create_card(slide17, left, top, Inches(5.6), Inches(2.2), title=wtitle, border_color=wcolor)
        tb = slide17.shapes.add_textbox(left + Inches(0.2), top + Inches(0.55), Inches(5.2), Inches(1.5))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = wdesc
        p.font.size = Pt(12)
        p.font.color.rgb = COLOR_TEXT_MAIN

    # ==================== SLIDE 18: COMPLETE DIGITAL PLATFORM ====================
    slide18 = prs.slides.add_slide(blank_layout)
    add_bg(slide18)
    add_header(slide18, "18", "COMPLETE DIGITAL PLATFORM MODULES", "Unified web application architecture for multi-stakeholder workflows")

    modules = [
        "Overview Dashboard", "Interactive GIS Map", "Live Sensor Telemetry",
        "Groundwater Potential", "Water Quality Index", "Recharge Mapping",
        "Historical Trends", "AI Stress Predictions", "Government Reports"
    ]

    for idx, mod in enumerate(modules):
        col = idx % 3
        row = idx // 3
        left = Inches(0.8 + col * 4.0)
        top = Inches(1.8 + row * 1.65)

        create_card(slide18, left, top, Inches(3.7), Inches(1.4), border_color=COLOR_CYAN)
        tb = slide18.shapes.add_textbox(left + Inches(0.2), top + Inches(0.2), Inches(3.3), Inches(1.0))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"MODULE 0{idx+1}"
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = COLOR_CYAN

        p2 = tf.add_paragraph()
        p2.text = mod
        p2.font.size = Pt(13)
        p2.font.bold = True
        p2.font.color.rgb = COLOR_TEXT_MAIN

    # ==================== SLIDE 19: TECHNOLOGY & AI STACK ====================
    slide19 = prs.slides.add_slide(blank_layout)
    add_bg(slide19)
    add_header(slide19, "19", "TECHNOLOGY & AI STACK", "End-to-end technology stack powering hardware, analytics, and visual narrative")

    stack_cats = [
        ("Hardware Layer", "Raspberry Pi 3B+, ESP32, 6 Physical Sensors, Custom PCB", COLOR_CYAN),
        ("Backend & Database", "Python Server, Relational/Spatial Database, REST Hardware API", COLOR_GREEN),
        ("GIS & Remote Sensing", "Geospatial Engine, Raster Processing, Vector Layer Indexing", COLOR_ACCENT_BLUE),
        ("AI / ML Engine", "Predictive Intelligence Layer, Feature Pipeline, Time-Series Modeling", COLOR_AMBER),
        ("NVIDIA AI Stack", "NVIDIA Nemotron-3-Ultra-550B-A55B (Reasoning) | FLUX.1-dev (Visuals)", COLOR_RED)
    ]

    for idx, (scat, sval, scolor) in enumerate(stack_cats):
        top = Inches(1.8 + idx * 1.0)
        create_card(slide19, Inches(0.8), top, Inches(11.7), Inches(0.85), border_color=scolor)
        tb = slide19.shapes.add_textbox(Inches(1.0), top + Inches(0.15), Inches(11.3), Inches(0.55))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"{scat.upper()}:  "
        p.font.bold = True
        p.font.size = Pt(12)
        p.font.color.rgb = scolor
        
        p2 = tf.add_paragraph()
        p2.text = f"   {sval}"
        p2.font.size = Pt(11)
        p2.font.color.rgb = COLOR_TEXT_MAIN

    # ==================== SLIDE 20: FUTURE SCOPE & CLOSING ====================
    slide20 = prs.slides.add_slide(blank_layout)
    add_bg(slide20)
    add_header(slide20, "20", "FUTURE SCOPE & ROADMAP", "From physical prototype to regional groundwater intelligence")

    steps_scope = [
        ("Phase 1: Prototype Node", "Single physical IoT node & core web dashboard validation."),
        ("Phase 2: Field Sensor Net", "Deploy multi-node sensor array across target watershed."),
        ("Phase 3: Multi-Location", "Scale edge hardware gateways across diverse hydrogeological zones."),
        ("Phase 4: Remote Sensing Fusion", "Automate high-resolution satellite raster ingestion."),
        ("Phase 5: AI Prediction", "Train & validate predictive models on multi-year temporal dataset."),
        ("Phase 6: Regional Intelligence", "Full government decision-support deployment for basin management.")
    ]

    for idx, (pstep, pdesc) in enumerate(steps_scope):
        col = idx % 2
        row = idx // 2
        left = Inches(0.8 + col * 5.9)
        top = Inches(1.8 + row * 1.6)

        create_card(slide20, left, top, Inches(5.6), Inches(1.35), border_color=COLOR_CYAN)
        tb = slide20.shapes.add_textbox(left + Inches(0.2), top + Inches(0.15), Inches(5.2), Inches(1.0))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = pstep
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = COLOR_CYAN

        p2 = tf.add_paragraph()
        p2.text = pdesc
        p2.font.size = Pt(10)
        p2.font.color.rgb = COLOR_TEXT_MUTED

    # Bottom Closing Statement
    create_card(slide20, Inches(0.8), Inches(5.8), Inches(11.7), Inches(1.1), title="CLOSING STATEMENT", border_color=COLOR_GREEN)
    tb_close = slide20.shapes.add_textbox(Inches(1.0), Inches(6.15), Inches(11.3), Inches(0.6))
    tf_close = tb_close.text_frame
    tf_close.word_wrap = True
    p = tf_close.paragraphs[0]
    p.text = '"FROM DATA TO GROUNDWATER INTELLIGENCE."'
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLOR_CYAN

    out_file = r"d:\Project Zero\docs\groundwater_intelligence_platform.pptx"
    os.makedirs(os.path.dirname(out_file), exist_ok=True)
    prs.save(out_file)
    print(f"SUCCESS: PowerPoint presentation created at {out_file}")

if __name__ == "__main__":
    build_presentation()
