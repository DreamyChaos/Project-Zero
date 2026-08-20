import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_sih_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Editorial Color System
    COLOR_BG = RGBColor(7, 11, 18)            # #070B12 Deep Space Black
    COLOR_CARD = RGBColor(15, 23, 42)         # #0F172A Subsurface Navy
    COLOR_CARD_LIGHT = RGBColor(30, 41, 59)   # #1E293B Card Light
    COLOR_BORDER = RGBColor(51, 65, 85)       # #334155 Slate Border
    COLOR_CYAN = RGBColor(0, 240, 255)         # #00F0FF Electric Cyan
    COLOR_BLUE_ACCENT = RGBColor(2, 132, 199)  # #0284C7 Water Blue
    COLOR_GREEN = RGBColor(16, 185, 129)      # #10B981 Hydro Green
    COLOR_TEAL = RGBColor(20, 184, 166)        # #14B8A6 Teal
    COLOR_PURPLE = RGBColor(124, 58, 237)      # #7C3AED Deep Purple Accent
    COLOR_AMBER = RGBColor(245, 158, 11)      # #F59E0B Warning Amber
    COLOR_CORAL = RGBColor(249, 115, 22)      # #F97316 Coral
    COLOR_RED = RGBColor(239, 68, 68)         # #EF4444 Alert Red
    COLOR_TEXT_MAIN = RGBColor(248, 250, 252) # #F8FAFC
    COLOR_TEXT_MUTED = RGBColor(148, 163, 184) # #94A3B8

    ASSETS_DIR = r"D:\SIH PPT\assets"
    COVER_IMG = os.path.join(ASSETS_DIR, "aquifer_cover.jpg")
    GIS_IMG = os.path.join(ASSETS_DIR, "satellite_gis.jpg")
    SPLIT_IMG = os.path.join(ASSETS_DIR, "surface_subsurface_problem.jpg")
    DASHBOARD_IMG = os.path.join(ASSETS_DIR, "dashboard_ui_product.jpg")

    def add_bg(slide):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = COLOR_BG
        bg.line.fill.background()
        return bg

    def add_header(slide, num_str, title_str, subtitle_str=""):
        tx_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.35))
        tf = tx_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"{num_str}  // GROUNDWATER INTELLIGENCE PLATFORM"
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = COLOR_CYAN
        p.font.name = 'Segoe UI'

        tx_box2 = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.75))
        tf2 = tx_box2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = title_str
        p2.font.size = Pt(22)
        p2.font.bold = True
        p2.font.color.rgb = COLOR_TEXT_MAIN
        p2.font.name = 'Segoe UI'

        if subtitle_str:
            p3 = tf2.add_paragraph()
            p3.text = subtitle_str
            p3.font.size = Pt(12)
            p3.font.color.rgb = COLOR_TEXT_MUTED
            p3.font.name = 'Segoe UI'

    def create_card(slide, left, top, width, height, title="", border_color=COLOR_BORDER, fill_color=COLOR_CARD):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = fill_color
        card.line.color.rgb = border_color
        card.line.width = Pt(1.5)
        
        if title:
            tb = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.15), width - Inches(0.4), Inches(0.4))
            tf = tb.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = title.upper()
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = COLOR_CYAN
            p.font.name = 'Segoe UI'
        return card

    # ==================== SLIDE 01: COVER ====================
    slide1 = prs.slides.add_slide(blank_layout)
    add_bg(slide1)
    if os.path.exists(COVER_IMG):
        slide1.shapes.add_picture(COVER_IMG, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    
    # Dark glassmorphic overlay box on left
    card1 = create_card(slide1, Inches(0.8), Inches(1.0), Inches(6.8), Inches(5.5), border_color=COLOR_CYAN)
    tb = slide1.shapes.add_textbox(Inches(1.1), Inches(1.3), Inches(6.2), Inches(5.0))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "GROUNDWATER\nINTELLIGENCE"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT_MAIN

    p2 = tf.add_paragraph()
    p2.text = "PLATFORM"
    p2.font.size = Pt(22)
    p2.font.bold = True
    p2.font.color.rgb = COLOR_CYAN

    p3 = tf.add_paragraph()
    p3.text = "\nA Data-Driven System for Monitoring, Mapping & Predictive Decision Support"
    p3.font.size = Pt(14)
    p3.font.color.rgb = COLOR_TEXT_MUTED

    p4 = tf.add_paragraph()
    p4.text = "\nSENSE → INTEGRATE → MAP → ANALYZE → PREDICT → DECIDE"
    p4.font.size = Pt(11)
    p4.font.bold = True
    p4.font.color.rgb = COLOR_GREEN

    p5 = tf.add_paragraph()
    p5.text = "\nReasoning: NVIDIA Nemotron-3-Ultra-550B-A55B\nVisuals: FLUX.1-dev via NVIDIA NIM"
    p5.font.size = Pt(10)
    p5.font.color.rgb = COLOR_TEXT_MUTED

    # ==================== SLIDE 02: THE PROBLEM ====================
    slide2 = prs.slides.add_slide(blank_layout)
    add_bg(slide2)
    add_header(slide2, "02", "GROUNDWATER IS INVISIBLE. ITS IMPACT ISN'T.", "Surface activity vs. subterranean reserve depletion")

    # Left: Headline + key statistics
    tb_left = slide2.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.6), Inches(5.0))
    tf_l = tb_left.text_frame
    tf_l.word_wrap = True
    
    p = tf_l.paragraphs[0]
    p.text = "THE UNSEEN CRITICAL CRISIS"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLOR_CORAL

    bullets = [
        ("Subterranean Depletion", "Unmonitored over-pumping causes severe water table drawdown and storage loss."),
        ("Water Quality Degradation", "Silent salinization, EC spikes, and pH contamination threaten public health."),
        ("Policy Lag & Blind Decisions", "Interventions arrive after wells run dry due to lack of real-time spatial intelligence.")
    ]
    for btitle, bdesc in bullets:
        p = tf_l.add_paragraph()
        p.text = f"\n• {btitle}: "
        p.font.bold = True
        p.font.size = Pt(12)
        p.font.color.rgb = COLOR_TEXT_MAIN
        p_sub = tf_l.add_paragraph()
        p_sub.text = f"  {bdesc}"
        p_sub.font.size = Pt(11)
        p_sub.font.color.rgb = COLOR_TEXT_MUTED

    # Right: Photorealistic split image
    if os.path.exists(SPLIT_IMG):
        slide2.shapes.add_picture(SPLIT_IMG, Inches(6.8), Inches(1.8), Inches(5.7), Inches(5.1))

    # ==================== SLIDE 03: THE DATA GAP ====================
    slide3 = prs.slides.add_slide(blank_layout)
    add_bg(slide3)
    add_header(slide3, "03", "THE DATA EXISTS. THE CONNECTION DOESN'T.", "Groundwater data is trapped in isolated, non-interoperable silos")

    silos = [
        ("01 / SENSOR LOGS", "Local well depth & water metrics", COLOR_CYAN),
        ("02 / SATELLITE RASTERS", "Soil moisture & GRACE data", COLOR_BLUE_ACCENT),
        ("03 / RAINFALL DATA", "Meteorological precipitation", COLOR_GREEN),
        ("04 / SOIL & GIS", "Infiltration & slope rasters", COLOR_AMBER),
        ("05 / LEGACY REPORTS", "State paper records & surveys", COLOR_CORAL)
    ]
    for idx, (stitle, sdesc, scolor) in enumerate(silos):
        top_pos = Inches(1.8 + idx * 1.0)
        create_card(slide3, Inches(0.8), top_pos, Inches(5.6), Inches(0.85), border_color=scolor)
        tb = slide3.shapes.add_textbox(Inches(1.0), top_pos + Inches(0.12), Inches(5.2), Inches(0.6))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"{stitle}  [ ISOLATED SILO ]"
        p.font.bold = True
        p.font.size = Pt(11)
        p.font.color.rgb = scolor
        p2 = tf.add_paragraph()
        p2.text = sdesc
        p2.font.size = Pt(10)
        p2.font.color.rgb = COLOR_TEXT_MUTED

    # Right Bottleneck Box
    create_card(slide3, Inches(6.8), Inches(1.8), Inches(5.7), Inches(5.1), title="THE FRAGMENTATION BOTTLENECK", border_color=COLOR_RED)
    tb_bot = slide3.shapes.add_textbox(Inches(7.1), Inches(2.5), Inches(5.1), Inches(4.0))
    tf_bot = tb_bot.text_frame
    tf_bot.word_wrap = True
    p = tf_bot.paragraphs[0]
    p.text = "Why Current Management Fails:"
    p.font.bold = True
    p.font.size = Pt(14)
    p.font.color.rgb = COLOR_TEXT_MAIN
    
    reasons = [
        "No single platform unifies physical sensors with satellite imagery.",
        "GIS layers remain static maps rather than dynamic predictive models.",
        "Decision-makers receive delayed reports after depletion occurs."
    ]
    for r in reasons:
        p = tf_bot.add_paragraph()
        p.text = f"\n❌ {r}"
        p.font.size = Pt(11)
        p.font.color.rgb = COLOR_TEXT_MUTED

    # ==================== SLIDE 04: OUR VISION ====================
    slide4 = prs.slides.add_slide(blank_layout)
    add_bg(slide4)
    add_header(slide4, "04", "ONE PLATFORM. MULTIPLE DATA SOURCES. ONE GROUNDWATER VIEW.", "Fusing physical telemetry with spatial environmental intelligence")

    create_card(slide4, Inches(0.8), Inches(1.8), Inches(5.6), Inches(5.1), title="DATA CONVERGENCE PIPELINE", border_color=COLOR_CYAN)
    tb = slide4.shapes.add_textbox(Inches(1.0), Inches(2.3), Inches(5.2), Inches(4.4))
    tf = tb.text_frame
    tf.word_wrap = True
    bullets = [
        ("Sensors + Satellite + GIS", "Fusing ground-truth telemetry with macro satellite rasters."),
        ("Centralized Processing", "Single source of truth feeding spatial analytics."),
        ("Continuous Ingestion", "Automated pipelines combining real-time and historical datasets.")
    ]
    for btitle, bdesc in bullets:
        p = tf.add_paragraph()
        p.text = f"• {btitle}: "
        p.font.bold = True
        p.font.size = Pt(12)
        p.font.color.rgb = COLOR_CYAN
        p_sub = tf.add_paragraph()
        p_sub.text = f"  {bdesc}\n"
        p_sub.font.size = Pt(11)
        p_sub.font.color.rgb = COLOR_TEXT_MUTED

    create_card(slide4, Inches(6.8), Inches(1.8), Inches(5.7), Inches(5.1), title="THE VALUE NARRATIVE", border_color=COLOR_GREEN)
    tb_right = slide4.shapes.add_textbox(Inches(7.1), Inches(2.3), Inches(5.1), Inches(4.4))
    tf_r = tb_right.text_frame
    tf_r.word_wrap = True
    steps = [
        "1. SENSE  → Real-time physical sensor sampling",
        "2. INTEGRATE → Multi-source data pipeline fusion",
        "3. MAP → GIS spatial potential & risk mapping",
        "4. ANALYZE → Time-series drawdown evaluation",
        "5. PREDICT → AI/ML predictive forecasting",
        "6. DECIDE → Government decision-support dashboard"
    ]
    for s in steps:
        p = tf_r.add_paragraph()
        p.text = s
        p.font.bold = True
        p.font.size = Pt(12)
        p.font.color.rgb = COLOR_TEXT_MAIN
        tf_r.add_paragraph().text = ""

    # ==================== SLIDE 05: SYSTEM ARCHITECTURE ====================
    slide5 = prs.slides.add_slide(blank_layout)
    add_bg(slide5)
    add_header(slide5, "05", "SYSTEM ARCHITECTURE", "End-to-end multi-layer software & data architecture pipeline")

    layers = [
        ("01 / PHYSICAL LAYER", "Raspberry Pi 3B+ | ESP32 | 6 Telemetry Sensors | Custom PCB", COLOR_CYAN),
        ("02 / GATEWAY & API", "Hardware Data API | Data Ingestion | Validation & Normalization", COLOR_BLUE_ACCENT),
        ("03 / DATABASE & GIS", "Spatial Database | Layer Indexing | Satellite Raster Processing", COLOR_GREEN),
        ("04 / ANALYTICS ENGINE", "Groundwater Potential Mapping | Recharge Assessment | Risk Scores", COLOR_AMBER),
        ("05 / PREDICTIVE AI/ML", "AI/ML Prediction Engine | Trend Analytics | Scenario Modeling", COLOR_RED),
        ("06 / DECISION PLATFORM", "Web Frontend | GIS Dashboard | Executive Decision Support Reports", COLOR_PURPLE)
    ]
    for idx, (ltitle, ldesc, lcolor) in enumerate(layers):
        top_pos = Inches(1.8 + idx * 0.88)
        create_card(slide5, Inches(0.8), top_pos, Inches(11.7), Inches(0.78), border_color=lcolor)
        tb = slide5.shapes.add_textbox(Inches(1.0), top_pos + Inches(0.12), Inches(11.3), Inches(0.55))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"{ltitle}  "
        p.font.bold = True
        p.font.size = Pt(11)
        p.font.color.rgb = lcolor
        p2 = tf.add_paragraph()
        p2.text = f"    {ldesc}"
        p2.font.size = Pt(10)
        p2.font.color.rgb = COLOR_TEXT_MUTED

    # ==================== SLIDE 06: FIELD DATA NODE ====================
    slide6 = prs.slides.add_slide(blank_layout)
    add_bg(slide6)
    add_header(slide6, "06", "FIELD DATA NODE", "Hardware edge gateway & ground-truth acquisition setup")

    create_card(slide6, Inches(0.8), Inches(1.8), Inches(5.6), Inches(5.1), title="HARDWARE CORE COMPONENTS", border_color=COLOR_CYAN)
    tb_hw = slide6.shapes.add_textbox(Inches(1.0), Inches(2.3), Inches(5.2), Inches(4.4))
    tf_hw = tb_hw.text_frame
    tf_hw.word_wrap = True
    hw_list = [
        ("Raspberry Pi 3B+", "Central edge gateway, local payload buffering & API dispatch."),
        ("ESP32 Microcontroller", "Low-latency sensor analog sampling & digital bus controller."),
        ("Custom PCB", "Regulated power delivery, signal filtering & connector bus."),
        ("Pump + Tubing", "Controlled fluid circulation for dynamic sensor testing."),
        ("Physical Model", "Aquifer test column for continuous prototype validation.")
    ]
    for htitle, hdesc in hw_list:
        p = tf_hw.add_paragraph()
        p.text = f"• {htitle}: "
        p.font.bold = True
        p.font.size = Pt(11)
        p.font.color.rgb = COLOR_TEXT_MAIN
        p_sub = tf_hw.add_paragraph()
        p_sub.text = f"  {hdesc}"
        p_sub.font.size = Pt(10)
        p_sub.font.color.rgb = COLOR_TEXT_MUTED

    create_card(slide6, Inches(6.8), Inches(1.8), Inches(5.7), Inches(5.1), title="HARDWARE POSITIONING & ROLE", border_color=COLOR_GREEN)
    tb_role = slide6.shapes.add_textbox(Inches(7.1), Inches(2.3), Inches(5.1), Inches(4.4))
    tf_role = tb_role.text_frame
    tf_role.word_wrap = True
    p = tf_role.paragraphs[0]
    p.text = "Hardware vs. Software Focus:"
    p.font.bold = True
    p.font.size = Pt(13)
    p.font.color.rgb = COLOR_CYAN

    roles = [
        ("Empirical Ground-Truth Provider", "Provides physical validation for spatial models & satellite rasters."),
        ("Continuous Edge Sampling", "Autonomous monitoring of water depth, pH, EC, and turbidity."),
        ("Software-Centric Focus", "The hardware is the edge data collector; the software platform is the primary project engine.")
    ]
    for rtitle, rdesc in roles:
        p = tf_role.add_paragraph()
        p.text = f"\n✓ {rtitle}"
        p.font.bold = True
        p.font.size = Pt(11)
        p.font.color.rgb = COLOR_GREEN
        p2 = tf_role.add_paragraph()
        p2.text = f"   {rdesc}"
        p2.font.size = Pt(10)
        p2.font.color.rgb = COLOR_TEXT_MUTED

    # ==================== SLIDE 07: SENSOR DATA ====================
    slide7 = prs.slides.add_slide(blank_layout)
    add_bg(slide7)
    add_header(slide7, "07", "SENSOR INSTRUMENTATION & TELEMETRY", "Key physical parameters captured by the edge payload (Illustrative Telemetry)")

    sensors = [
        ("Groundwater Level", "Ultrasonic / Depth Probe", "42.5 m depth", COLOR_CYAN),
        ("pH Sensor", "Electrochemical Electrode", "7.2 pH (Optimal)", COLOR_GREEN),
        ("Electrical Conductivity", "EC Conductivity Probe", "450 µS/cm", COLOR_AMBER),
        ("Turbidity Sensor", "Optical Scattering Probe", "3.1 NTU", COLOR_RED),
        ("Temperature Sensor", "Digital Thermal Probe", "22.4 °C", COLOR_BLUE_ACCENT),
        ("Soil-Moisture Sensor", "Capacitive Moisture Sensor", "34% Saturation", COLOR_PURPLE)
    ]

    for idx, (stitle, ssub, sval, scolor) in enumerate(sensors):
        col = idx % 3
        row = idx // 3
        left_pos = Inches(0.8 + col * 4.0)
        top_pos = Inches(1.8 + row * 2.5)

        create_card(slide7, left_pos, top_pos, Inches(3.7), Inches(2.2), title=stitle, border_color=scolor)
        tb = slide7.shapes.add_textbox(left_pos + Inches(0.15), top_pos + Inches(0.5), Inches(3.4), Inches(1.6))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].text = ssub
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.size = Pt(10)
        tf.paragraphs[0].font.color.rgb = COLOR_TEXT_MUTED

        p = tf.add_paragraph()
        p.text = f"\nSample Value: {sval}"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = scolor

        p2 = tf.add_paragraph()
        p2.text = "[ ILLUSTRATIVE TELEMETRY METRIC ]"
        p2.font.size = Pt(9)
        p2.font.color.rgb = COLOR_TEXT_MUTED

    # ==================== SLIDE 08: GIS INTERACTIVE MAP ====================
    slide8 = prs.slides.add_slide(blank_layout)
    add_bg(slide8)
    add_header(slide8, "08", "TURNING DATA INTO A LIVING MAP", "GIS interactive spatial intelligence & multi-layer visualization engine")

    create_card(slide8, Inches(0.8), Inches(1.8), Inches(5.6), Inches(5.1), title="SPATIAL MAPPING CAPABILITIES", border_color=COLOR_CYAN)
    tb = slide8.shapes.add_textbox(Inches(1.0), Inches(2.3), Inches(5.2), Inches(4.4))
    tf = tb.text_frame
    tf.word_wrap = True
    gis_features = [
        ("Multi-Layer Geospatial Overlay", "Integrate well telemetry, administrative borders, topography, and watershed basins."),
        ("Interactive Spatial Queries", "Clickable well nodes displaying live depth history, pH trends, and risk scores."),
        ("Raster & Vector Processing", "Smooth rendering of satellite rasters alongside point vector data."),
        ("Spatial Heatmap Interpolation", "Estimate groundwater potential across unmonitored zones between physical wells.")
    ]
    for ftitle, fdesc in gis_features:
        p = tf.add_paragraph()
        p.text = f"• {ftitle}: "
        p.font.bold = True
        p.font.size = Pt(11)
        p.font.color.rgb = COLOR_CYAN
        p_sub = tf.add_paragraph()
        p_sub.text = f"  {fdesc}\n"
        p_sub.font.size = Pt(10)
        p_sub.font.color.rgb = COLOR_TEXT_MUTED

    if os.path.exists(GIS_IMG):
        slide8.shapes.add_picture(GIS_IMG, Inches(6.8), Inches(1.8), Inches(5.7), Inches(5.1))

    # ==================== SLIDE 09: REMOTE SENSING ====================
    slide9 = prs.slides.add_slide(blank_layout)
    add_bg(slide9)
    add_header(slide9, "09", "SEEING WHAT THE GROUND CAN'T TELL US DIRECTLY", "Satellite remote sensing & macro environmental data fusion")

    cards_rs = [
        ("01 / Satellite Imagery", "Multi-spectral optical & radar data for surface soil moisture and vegetation health (NDVI).", COLOR_CYAN),
        ("02 / Environmental Logs", "Continuous ingestion of precipitation, temperature, evapotranspiration, and humidity datasets.", COLOR_GREEN),
        ("03 / Spatial Gap Filling", "Fills spatial gaps between physical well sensors without requiring thousands of physical nodes.", COLOR_AMBER)
    ]
    for idx, (title, desc, color) in enumerate(cards_rs):
        left = Inches(0.8 + idx * 4.0)
        create_card(slide9, left, Inches(1.8), Inches(3.7), Inches(5.1), title=title, border_color=color)
        tb = slide9.shapes.add_textbox(left + Inches(0.2), Inches(2.4), Inches(3.3), Inches(4.2))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(11)
        p.font.color.rgb = COLOR_TEXT_MAIN

    # ==================== SLIDE 10: GROUNDWATER POTENTIAL MAPPING ====================
    slide10 = prs.slides.add_slide(blank_layout)
    add_bg(slide10)
    add_header(slide10, "10", "GROUNDWATER POTENTIAL MAPPING", "Multi-criteria spatial assessment for estimated groundwater availability")

    create_card(slide10, Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.1), title="SPATIAL WEIGHTED OVERLAY MATRIX", border_color=COLOR_CYAN)
    tb = slide10.shapes.add_textbox(Inches(1.1), Inches(2.3), Inches(11.1), Inches(4.4))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Assessment Model: Potential Score = f(Rainfall, Soil, Slope, Land Cover, Drainage, Telemetry)"
    p.font.bold = True
    p.font.size = Pt(13)
    p.font.color.rgb = COLOR_TEXT_MAIN

    factors = [
        ("Topography & Slope", "Flatter terrain promotes water ponding and surface infiltration over rapid runoff."),
        ("Soil Permeability", "Porous sandy soils enable higher percolation rates than dense clay layers."),
        ("Drainage Density", "Lower stream density indicates higher subsurface absorption capacity."),
        ("Land Cover & NDVI", "Dense vegetation enhances root-zone soil porosity and reduces direct evaporation.")
    ]
    for ftitle, fdesc in factors:
        p = tf.add_paragraph()
        p.text = f"\n• {ftitle}: "
        p.font.bold = True
        p.font.color.rgb = COLOR_CYAN
        p.font.size = Pt(11)
        p2 = tf.add_paragraph()
        p2.text = f"  {fdesc}"
        p2.font.size = Pt(10)
        p2.font.color.rgb = COLOR_TEXT_MUTED

    # ==================== SLIDE 11: RECHARGE-POTENTIAL MAPPING ====================
    slide11 = prs.slides.add_slide(blank_layout)
    add_bg(slide11)
    add_header(slide11, "11", "RECHARGE-POTENTIAL MAPPING", "Identifying optimal zones for artificial & natural groundwater recharge")

    col_w = Inches(5.6)
    create_card(slide11, Inches(0.8), Inches(1.8), col_w, Inches(5.1), title="INFILTRATION & RECHARGE FACTORS", border_color=COLOR_GREEN)
    tb1 = slide11.shapes.add_textbox(Inches(1.0), Inches(2.3), col_w - Inches(0.4), Inches(4.4))
    tf1 = tb1.text_frame
    tf1.word_wrap = True
    r_factors = [
        "Precipitation Intensity & Duration",
        "Soil Hydraulic Conductivity",
        "Geological Lithology & Permeability",
        "Surface Slope & Runoff Velocity",
        "Soil Moisture Saturation Baseline"
    ]
    for rf in r_factors:
        p = tf1.add_paragraph()
        p.text = f"• {rf}"
        p.font.size = Pt(11)
        p.font.color.rgb = COLOR_TEXT_MUTED

    create_card(slide11, Inches(6.8), Inches(1.8), col_w, Inches(5.1), title="ACTIONABLE DECISION INTERVENTIONS", border_color=COLOR_CYAN)
    tb2 = slide11.shapes.add_textbox(Inches(7.0), Inches(2.3), col_w - Inches(0.4), Inches(4.4))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    actions = [
        ("Check Dam Site Selection", "Identify high-percolation zones for check dams."),
        ("Recharge Shaft Placement", "Locate optimal coordinates for injection shafts."),
        ("Rainwater Harvesting Strategy", "Guide municipal rainwater collection mandates."),
        ("Conservation Zoning", "Designate protected high-recharge watershed basins.")
    ]
    for atitle, adesc in actions:
        p = tf2.add_paragraph()
        p.text = f"✓ {atitle}: "
        p.font.bold = True
        p.font.size = Pt(11)
        p.font.color.rgb = COLOR_GREEN
        p_sub = tf2.add_paragraph()
        p_sub.text = f"  {adesc}"
        p_sub.font.size = Pt(10)
        p_sub.font.color.rgb = COLOR_TEXT_MUTED

    # ==================== SLIDE 12: WATER-QUALITY RISK MAPPING ====================
    slide12 = prs.slides.add_slide(blank_layout)
    add_bg(slide12)
    add_header(slide12, "12", "WATER-QUALITY RISK MAPPING", "Transforming well sensor measurements into spatial risk maps")

    q_cards = [
        ("pH Spatial Index", "Acidity & Alkalinity", "Maps agricultural runoff, chemical leaching & industrial contamination.", COLOR_CYAN),
        ("EC Risk Index", "Electrical Conductivity", "Tracks salinization, mineral dissolving & coastal saltwater intrusion.", COLOR_AMBER),
        ("Turbidity Risk Index", "Suspended Sediment", "Identifies wellhead structural damage, sediment ingress & filtration failure.", COLOR_RED)
    ]
    for idx, (qtitle, qsub, qdesc, qcolor) in enumerate(q_cards):
        left = Inches(0.8 + idx * 4.0)
        create_card(slide12, left, Inches(1.8), Inches(3.7), Inches(3.2), title=qtitle, border_color=qcolor)
        tb = slide12.shapes.add_textbox(left + Inches(0.2), Inches(2.4), Inches(3.3), Inches(2.4))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].text = qsub
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.size = Pt(11)
        tf.paragraphs[0].font.color.rgb = COLOR_TEXT_MAIN
        p = tf.add_paragraph()
        p.text = f"\n{qdesc}"
        p.font.size = Pt(10)
        p.font.color.rgb = COLOR_TEXT_MUTED

    create_card(slide12, Inches(0.8), Inches(5.3), Inches(11.7), Inches(1.6), title="INTEGRATED WATER-QUALITY RISK SCORE", border_color=COLOR_GREEN)
    tb_bot = slide12.shapes.add_textbox(Inches(1.0), Inches(5.65), Inches(11.3), Inches(1.1))
    tf_bot = tb_bot.text_frame
    tf_bot.word_wrap = True
    p = tf_bot.paragraphs[0]
    p.text = "Spatial interpolation combines pH, EC, and Turbidity into a unified Risk Index (Safe / Caution / Hazardous) to alert authorities prior to public water supply contamination."
    p.font.size = Pt(12)
    p.font.color.rgb = COLOR_TEXT_MAIN

    # ==================== SLIDE 13: HISTORICAL TREND ANALYSIS ====================
    slide13 = prs.slides.add_slide(blank_layout)
    add_bg(slide13)
    add_header(slide13, "13", "HISTORICAL TREND ANALYSIS", "Time-series visualization of water table depth and environmental factors (Illustrative Data)")

    create_card(slide13, Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.1), title="TEMPORAL PATTERN EVALUATION", border_color=COLOR_CYAN)
    tb = slide13.shapes.add_textbox(Inches(1.1), Inches(2.3), Inches(11.1), Inches(4.4))
    tf = tb.text_frame
    tf.word_wrap = True
    trends = [
        ("Seasonal Water-Table Fluctuations", "Comparing monsoon recharge peaks against summer drawdown troughs."),
        ("Multi-Year Depletion Trajectories", "Identifying long-term negative baseline drift across consecutive years."),
        ("Rainfall-Recharge Lag Analysis", "Quantifying response time lag between rainfall events and water table rise."),
        ("Anomalous Event Detection", "Flagging sudden unseasonal drops caused by over-pumping or sensor faults.")
    ]
    for ttitle, tdesc in trends:
        p = tf.add_paragraph()
        p.text = f"\n• {ttitle}: "
        p.font.bold = True
        p.font.color.rgb = COLOR_CYAN
        p.font.size = Pt(11)
        p_sub = tf.add_paragraph()
        p_sub.text = f"  {tdesc}"
        p_sub.font.size = Pt(10)
        p_sub.font.color.rgb = COLOR_TEXT_MUTED

    # ==================== SLIDE 14: GROUNDWATER DEPLETION / STRESS ====================
    slide14 = prs.slides.add_slide(blank_layout)
    add_bg(slide14)
    add_header(slide14, "14", "GROUNDWATER DEPLETION & STRESS ZONES", "Identifying critical zones approaching threshold depletion")

    create_card(slide14, Inches(0.8), Inches(1.8), Inches(5.6), Inches(5.1), title="STRESS CLASSIFICATION MATRIX", border_color=COLOR_AMBER)
    tb1 = slide14.shapes.add_textbox(Inches(1.0), Inches(2.3), Inches(5.2), Inches(4.4))
    tf1 = tb1.text_frame
    tf1.word_wrap = True
    levels = [
        ("Safe Zone", "Stable water table, positive annual recharge balance.", COLOR_GREEN),
        ("Semi-Critical Zone", "Minor drawdown trend, extraction equals natural recharge.", COLOR_AMBER),
        ("Critical Stress Zone", "Severe persistent drawdown, extraction exceeds recharge by >20%.", COLOR_RED),
        ("Over-Exploited Zone", "Irreversible storage decline, risk of well failure & land subsidence.", COLOR_RED)
    ]
    for ltitle, ldesc, lcolor in levels:
        p = tf1.add_paragraph()
        p.text = f"■ {ltitle}: "
        p.font.bold = True
        p.font.color.rgb = lcolor
        p.font.size = Pt(11)
        p_sub = tf1.add_paragraph()
        p_sub.text = f"  {ldesc}\n"
        p_sub.font.size = Pt(10)
        p_sub.font.color.rgb = COLOR_TEXT_MUTED

    create_card(slide14, Inches(6.8), Inches(1.8), Inches(5.7), Inches(5.1), title="EARLY WARNING METRICS", border_color=COLOR_RED)
    tb2 = slide14.shapes.add_textbox(Inches(7.0), Inches(2.3), Inches(5.3), Inches(4.4))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    metrics = [
        ("Rate of Drawdown (m/year)", "Accelerating annual depletion rate."),
        ("Recharge Deficit Score", "Percentage of unreplenished extracted volume."),
        ("Water Quality Degradation Rate", "Simultaneous rise in EC/salinity alongside depth drop."),
        ("Proximity to Critical Infrastructure", "Threat level to municipal supply wells.")
    ]
    for mtitle, mdesc in metrics:
        p = tf2.add_paragraph()
        p.text = f"⚠ {mtitle}: "
        p.font.bold = True
        p.font.color.rgb = COLOR_TEXT_MAIN
        p.font.size = Pt(11)
        p_sub = tf2.add_paragraph()
        p_sub.text = f"  {mdesc}\n"
        p_sub.font.size = Pt(10)
        p_sub.font.color.rgb = COLOR_TEXT_MUTED

    # ==================== SLIDE 15: AI / ML ENGINE ====================
    slide15 = prs.slides.add_slide(blank_layout)
    add_bg(slide15)
    add_header(slide15, "15", "AI / ML PREDICTION ENGINE", "Predictive Intelligence Layer for groundwater forecasting (Algorithm Agnostic)")

    create_card(slide15, Inches(0.8), Inches(1.8), Inches(11.7), Inches(1.1), title="PREDICTIVE INTELLIGENCE ARCHITECTURE", border_color=COLOR_CYAN)
    tb_top = slide15.shapes.add_textbox(Inches(1.0), Inches(2.15), Inches(11.3), Inches(0.7))
    tf_top = tb_top.text_frame
    tf_top.word_wrap = True
    p = tf_top.paragraphs[0]
    p.text = "Described generically as an AI/ML Prediction Engine. The specific algorithm will be selected later based on dataset spatial-temporal resolution, feature quality, and empirical validation metrics."
    p.font.size = Pt(11)
    p.font.color.rgb = COLOR_TEXT_MAIN

    create_card(slide15, Inches(0.8), Inches(3.1), Inches(5.6), Inches(3.8), title="FEATURE PIPELINE INPUTS", border_color=COLOR_BLUE_ACCENT)
    tb_in = slide15.shapes.add_textbox(Inches(1.0), Inches(3.6), Inches(5.2), Inches(3.1))
    tf_in = tb_in.text_frame
    tf_in.word_wrap = True
    in_features = [
        "Historical Water Level Time Series",
        "Precipitation & Evapotranspiration Rates",
        "Soil Moisture & Satellite NDVI Rasters",
        "Geological Lithology & Infiltration Rates",
        "Live Edge Sensor Telemetry Feeds"
    ]
    for inf in in_features:
        p = tf_in.add_paragraph()
        p.text = f"→ {inf}"
        p.font.size = Pt(10)
        p.font.color.rgb = COLOR_TEXT_MUTED

    create_card(slide15, Inches(6.8), Inches(3.1), Inches(5.7), Inches(3.8), title="PREDICTIVE OUTPUTS", border_color=COLOR_GREEN)
    tb_out = slide15.shapes.add_textbox(Inches(7.0), Inches(3.6), Inches(5.3), Inches(3.1))
    tf_out = tb_out.text_frame
    tf_out.word_wrap = True
    out_features = [
        "30 / 60 / 90-Day Water Level Forecasts",
        "Depletion Risk Probability Distribution",
        "Recharge Deficit Predictions",
        "Water Quality Risk Trend Horizon",
        "Uncertainty Bounds & Confidence Intervals"
    ]
    for outf in out_features:
        p = tf_out.add_paragraph()
        p.text = f"★ {outf}"
        p.font.size = Pt(10)
        p.font.color.rgb = COLOR_TEXT_MAIN

    # ==================== SLIDE 16: FUTURE STRESS PREDICTION ====================
    slide16 = prs.slides.add_slide(blank_layout)
    add_bg(slide16)
    add_header(slide16, "16", "FUTURE GROUNDWATER-STRESS PREDICTION", "Scenario-based forecasting for climate variability and extraction demand")

    scenarios = [
        ("Baseline Scenario", "Normal Rainfall & Current Extraction", "Projects baseline water table behavior under standard seasonal precipitation.", COLOR_GREEN),
        ("Monsoon Deficit Scenario", "Below-Average Rainfall (-30%)", "Forecasts severe drawdown speed during drought years to prepare buffer storage.", COLOR_AMBER),
        ("High Extraction Demand", "Agricultural/Urban Expansion (+25%)", "Simulates aquifer stress under accelerated pumping schedules.", COLOR_RED)
    ]
    for idx, (stitle, ssub, sdesc, scolor) in enumerate(scenarios):
        left = Inches(0.8 + idx * 4.0)
        create_card(slide16, left, Inches(1.8), Inches(3.7), Inches(5.1), title=stitle, border_color=scolor)
        tb = slide16.shapes.add_textbox(left + Inches(0.2), Inches(2.4), Inches(3.3), Inches(4.2))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].text = ssub
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.size = Pt(11)
        tf.paragraphs[0].font.color.rgb = COLOR_TEXT_MAIN
        p = tf.add_paragraph()
        p.text = f"\n{sdesc}"
        p.font.size = Pt(10)
        p.font.color.rgb = COLOR_TEXT_MUTED

    # ==================== SLIDE 17: DECISION-SUPPORT DASHBOARD ====================
    slide17 = prs.slides.add_slide(blank_layout)
    add_bg(slide17)
    add_header(slide17, "17", "GOVERNMENT DECISION-SUPPORT DASHBOARD", "Turning complex spatial data and predictions into government policy action")

    if os.path.exists(DASHBOARD_IMG):
        slide17.shapes.add_picture(DASHBOARD_IMG, Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.1))

    # ==================== SLIDE 18: COMPLETE DIGITAL PLATFORM ====================
    slide18 = prs.slides.add_slide(blank_layout)
    add_bg(slide18)
    add_header(slide18, "18", "COMPLETE DIGITAL PLATFORM MODULE MATRIX", "Unified web application architecture for multi-stakeholder workflows")

    modules = [
        "Overview Dashboard", "Interactive GIS Map", "Live Sensor Telemetry",
        "Groundwater Potential", "Water Quality Index", "Recharge Mapping",
        "Historical Trends", "AI Stress Predictions", "Government Reports"
    ]
    for idx, mod in enumerate(modules):
        col = idx % 3
        row = idx // 3
        left = Inches(0.8 + col * 4.0)
        top = Inches(1.8 + row * 1.65)

        create_card(slide18, left, top, Inches(3.7), Inches(1.4), border_color=COLOR_CYAN)
        tb = slide18.shapes.add_textbox(left + Inches(0.2), top + Inches(0.2), Inches(3.3), Inches(1.0))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"MODULE 0{idx+1}"
        p.font.size = Pt(9)
        p.font.bold = True
        p.font.color.rgb = COLOR_CYAN

        p2 = tf.add_paragraph()
        p2.text = mod
        p2.font.size = Pt(12)
        p2.font.bold = True
        p2.font.color.rgb = COLOR_TEXT_MAIN

    # ==================== SLIDE 19: TECHNOLOGY & AI STACK ====================
    slide19 = prs.slides.add_slide(blank_layout)
    add_bg(slide19)
    add_header(slide19, "19", "TECHNOLOGY & AI MODEL STACK", "End-to-end technology stack powering hardware, analytics, and visual narrative")

    stack_cats = [
        ("PHYSICAL LAYER", "Raspberry Pi 3B+, ESP32 Microcontroller, 6 Telemetry Sensors, Custom PCB", COLOR_CYAN),
        ("DATA & BACKEND", "Python Server, Relational & Spatial Database, REST Hardware Data API, GIS Engine", COLOR_GREEN),
        ("PREDICTIVE ENGINE", "AI/ML Prediction Engine, Feature Pipeline, Time-Series Forecast Models", COLOR_AMBER),
        ("PRIMARY REASONING MODEL", "NVIDIA Nemotron-3-Ultra-550B-A55B (Content, Technical Writing & Narrative)", COLOR_RED),
        ("PRIMARY VISUAL MODEL", "FLUX.1-dev via NVIDIA NIM & NVIDIA MCP Integrations", COLOR_PURPLE)
    ]

    for idx, (scat, sval, scolor) in enumerate(stack_cats):
        top = Inches(1.8 + idx * 1.0)
        create_card(slide19, Inches(0.8), top, Inches(11.7), Inches(0.85), border_color=scolor)
        tb = slide19.shapes.add_textbox(Inches(1.0), top + Inches(0.15), Inches(11.3), Inches(0.55))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"{scat}:  "
        p.font.bold = True
        p.font.size = Pt(11)
        p.font.color.rgb = scolor
        p2 = tf.add_paragraph()
        p2.text = f"   {sval}"
        p2.font.size = Pt(11)
        p2.font.color.rgb = COLOR_TEXT_MAIN

    # ==================== SLIDE 20: FUTURE SCOPE & CLOSING ====================
    slide20 = prs.slides.add_slide(blank_layout)
    add_bg(slide20)
    add_header(slide20, "20", "FUTURE SCOPE & ROADMAP", "Progression roadmap from prototype node to regional groundwater intelligence")

    if os.path.exists(COVER_IMG):
        slide20.shapes.add_picture(COVER_IMG, Inches(0), Inches(0), Inches(13.333), Inches(7.5))

    card20 = create_card(slide20, Inches(0.8), Inches(1.2), Inches(11.7), Inches(5.3), border_color=COLOR_GREEN)
    tb20 = slide20.shapes.add_textbox(Inches(1.1), Inches(1.5), Inches(11.1), Inches(4.7))
    tf20 = tb20.text_frame
    tf20.word_wrap = True
    
    p = tf20.paragraphs[0]
    p.text = "FROM DATA TO GROUNDWATER INTELLIGENCE."
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLOR_CYAN

    p2 = tf20.add_paragraph()
    p2.text = "\nSENSE → INTEGRATE → MAP → ANALYZE → PREDICT → DECIDE"
    p2.font.size = Pt(16)
    p2.font.bold = True
    p2.font.color.rgb = COLOR_GREEN

    p3 = tf20.add_paragraph()
    p3.text = "\nScaling Roadmap:\nPrototype Node → Field Sensor Array → Multi-Location Gateways → Satellite Data Fusion → Predictive AI → Regional Intelligence"
    p3.font.size = Pt(13)
    p3.font.color.rgb = COLOR_TEXT_MAIN

    out_file = r"D:\SIH PPT\groundwater_intelligence_platform.pptx"
    prs.save(out_file)
    print(f"SUCCESS: Created redesigned presentation at {out_file}")

if __name__ == "__main__":
    create_sih_presentation()
